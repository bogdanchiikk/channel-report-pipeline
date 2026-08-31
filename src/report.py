import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from charts import plot_joins_unsubs, plot_subscribers
from metrics import load_metrics


def _add_title(slide, text, top=0.35):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True


def _add_body(slide, lines, top=1.3):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)


def build_report(sample_path: str, output_dir: str) -> Path:
    metrics = load_metrics(sample_path)
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    subs_png = plot_subscribers(metrics["daily"], out / "subscribers.png")
    joins_png = plot_joins_unsubs(metrics["daily"], out / "joins_unsubs.png")

    prs = Presentation()
    blank = prs.slide_layouts[6]

    title = prs.slides.add_slide(blank)
    _add_title(title, metrics["channel"] or "Channel report")
    _add_body(title, [f"Period: {metrics['period']}", "Source: sample JSON (not client data)"])

    kpi = prs.slides.add_slide(blank)
    _add_title(kpi, "Key numbers")
    _add_body(
        kpi,
        [
            f"Start: {metrics['start_subs']}",
            f"End: {metrics['end_subs']}",
            f"Net growth: {metrics['net_growth']}",
            f"Unsubs: {metrics['unsubs']}",
            f"Posts: {len(metrics['posts'])}",
        ],
    )

    for png, heading in (
        (subs_png, "Subscribers"),
        (joins_png, "Joined vs left"),
    ):
        slide = prs.slides.add_slide(blank)
        _add_title(slide, heading, top=0.2)
        slide.shapes.add_picture(str(png), Inches(0.5), Inches(1.1), width=Inches(9))

    top = prs.slides.add_slide(blank)
    _add_title(top, "Top posts by views")
    posts = metrics["posts"]
    if posts is None or posts.empty:
        _add_body(top, ["No posts"])
    else:
        ranked = posts.sort_values("views", ascending=False).head(5)
        lines = [
            f"{row['title']} — {int(row['views'])} views"
            for _, row in ranked.iterrows()
        ]
        _add_body(top, lines)

    pptx_path = out / "report.pptx"
    prs.save(pptx_path)
    return pptx_path


if __name__ == "__main__":
    config = json.loads(Path("config.example.json").read_text(encoding="utf-8"))
    path = build_report(config["sample"], config["output_dir"])
    print("saved", path)
